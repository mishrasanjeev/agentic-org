"""MIME-specific text extractors for the canonical RAG ingestion service.

Each extractor takes a ``bytes`` stream + filename and returns an
``ExtractedContent`` with the full text plus per-section provenance
(page / sheet / frame timestamp). Extractors that cannot produce usable
text raise ``UnsupportedMimeType`` — the ingestion service translates
that into a 415 at the API boundary rather than fake-indexing.

Native text is preferred. Scanned PDF pages and image documents use
Tesseract OCR with confidence and page provenance. Legacy/ODF Office files
are converted in an isolated temporary directory by headless LibreOffice.
Audio/video remain explicit unsupported types; uploads never report success
for a body that the platform cannot extract.
"""

from __future__ import annotations

import csv
import io
import json
import os
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from email import policy
from email.parser import BytesParser
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

MAX_PDF_PAGES = int(os.getenv("AGENTICORG_RAG_MAX_PDF_PAGES", "500"))
MAX_OCR_PAGES = int(os.getenv("AGENTICORG_RAG_MAX_OCR_PAGES", "100"))
MAX_IMAGE_PIXELS = int(os.getenv("AGENTICORG_RAG_MAX_IMAGE_PIXELS", "50000000"))
MAX_OFFICE_MEMBERS = int(os.getenv("AGENTICORG_RAG_MAX_OFFICE_MEMBERS", "10000"))
MAX_OFFICE_EXPANDED_BYTES = int(os.getenv("AGENTICORG_RAG_MAX_OFFICE_EXPANDED_BYTES", str(250 * 1024 * 1024)))


class UnsupportedMimeType(ValueError):  # noqa: N818 — external-facing API name; Error suffix would read redundantly
    """Raised by extractors + ingest service when a type isn't supported.

    Callers at the API boundary translate this into a ``415 Unsupported
    Media Type`` — distinct from a 422 validation error because the
    platform COULD accept the body, it just has no extractor wired.
    """


@dataclass
class ExtractedSpan:
    """A single chunk of extracted text with provenance."""

    text: str
    # Page number (1-indexed) for PDFs; None for other modalities.
    page: int | None = None
    # Sheet name for XLSX; None otherwise.
    sheet: str | None = None
    # Cell range for XLSX (e.g. "A1:Z42"); None otherwise.
    cell_range: str | None = None
    # Frame timestamp (seconds) for video; None otherwise.
    frame_timestamp_s: float | None = None


@dataclass
class ExtractedContent:
    """Result of extraction for a single uploaded artifact."""

    spans: list[ExtractedSpan]
    mime_type: str
    extraction_method: str  # "pypdf", "python-docx", "openpyxl", "text", ...
    total_chars: int = 0
    extra: dict[str, Any] = field(default_factory=dict)

    def full_text(self, separator: str = "\n\n") -> str:
        return separator.join(s.text for s in self.spans if s.text)


# ── Pure-text extractors ─────────────────────────────────────────────


_TEXT_LIKE_MIMETYPES = {
    "text/plain",
    "text/markdown",
    "text/csv",
    "application/json",
    "application/jsonl",
    "application/xml",
    "application/yaml",
    "application/x-yaml",
    "text/html",
    "application/xhtml+xml",
}


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self._hidden = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"script", "style", "noscript"}:
            self._hidden += 1
        elif tag in {"p", "br", "li", "tr", "h1", "h2", "h3", "h4"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript"} and self._hidden:
            self._hidden -= 1

    def handle_data(self, data: str) -> None:
        if not self._hidden and data.strip():
            self.parts.append(data.strip())


def _normalise_text(text: str) -> str:
    return "\n".join(line.strip() for line in text.splitlines() if line.strip())


def _validate_zip_container(stream: bytes) -> None:
    """Reject suspicious OOXML/ODF archives before parser expansion."""
    import zipfile

    try:
        with zipfile.ZipFile(io.BytesIO(stream)) as archive:
            members = archive.infolist()
            expanded = sum(member.file_size for member in members)
            if len(members) > MAX_OFFICE_MEMBERS or expanded > MAX_OFFICE_EXPANDED_BYTES:
                raise UnsupportedMimeType("Office archive exceeds safe expansion limits")
            for member in members:
                path = Path(member.filename)
                if path.is_absolute() or ".." in path.parts:
                    raise UnsupportedMimeType("Office archive contains an unsafe path")
    except zipfile.BadZipFile as exc:
        raise UnsupportedMimeType("Office document is not a valid ZIP container") from exc


def _extract_plaintext(stream: bytes, mime_type: str) -> ExtractedContent:
    try:
        text = stream.decode("utf-8")
    except UnicodeDecodeError:
        text = stream.decode("latin-1", errors="replace")
    return ExtractedContent(
        spans=[ExtractedSpan(text=text.strip())],
        mime_type=mime_type,
        extraction_method="text",
        total_chars=len(text),
    )


def _extract_html(stream: bytes, mime_type: str) -> ExtractedContent:
    parser = _HTMLTextExtractor()
    parser.feed(stream.decode("utf-8", errors="replace"))
    text = _normalise_text(" ".join(parser.parts))
    return ExtractedContent(
        spans=[ExtractedSpan(text=text)] if text else [],
        mime_type=mime_type,
        extraction_method="html-parser",
        total_chars=len(text),
    )


def _extract_eml(stream: bytes, mime_type: str) -> ExtractedContent:
    message = BytesParser(policy=policy.default).parsebytes(stream)
    parts = [
        f"Subject: {message.get('subject', '')}",
        f"From: {message.get('from', '')}",
        f"To: {message.get('to', '')}",
    ]
    if message.is_multipart():
        for part in message.walk():
            if part.get_content_type() == "text/plain" and not part.get_filename():
                parts.append(str(part.get_content()))
    elif message.get_content_type() == "text/plain":
        parts.append(str(message.get_content()))
    text = _normalise_text("\n".join(parts))
    return ExtractedContent(
        spans=[ExtractedSpan(text=text)] if text else [],
        mime_type=mime_type,
        extraction_method="email-parser",
        total_chars=len(text),
    )


def _extract_csv(stream: bytes, mime_type: str) -> ExtractedContent:
    """Flatten CSV rows into one chunk per row with row provenance."""
    try:
        text = stream.decode("utf-8-sig")
    except UnicodeDecodeError:
        text = stream.decode("latin-1", errors="replace")
    reader = csv.reader(io.StringIO(text))
    spans: list[ExtractedSpan] = []
    for idx, row in enumerate(reader):
        joined = " | ".join(cell.strip() for cell in row if cell.strip())
        if not joined:
            continue
        spans.append(ExtractedSpan(text=joined, cell_range=f"row {idx + 1}"))
    return ExtractedContent(
        spans=spans,
        mime_type=mime_type,
        extraction_method="csv",
        total_chars=sum(len(s.text) for s in spans),
    )


def _extract_json(stream: bytes, mime_type: str) -> ExtractedContent:
    try:
        text = stream.decode("utf-8")
        payload = json.loads(text)
        # Serialize back with indentation so retrieval can match on
        # structured keys.
        pretty = json.dumps(payload, indent=2, sort_keys=True, ensure_ascii=False)
        return ExtractedContent(
            spans=[ExtractedSpan(text=pretty)],
            mime_type=mime_type,
            extraction_method="json",
            total_chars=len(pretty),
        )
    except (UnicodeDecodeError, json.JSONDecodeError) as exc:
        raise UnsupportedMimeType(f"JSON body failed to decode for mime={mime_type!r}: {exc}") from exc


# ── PDF ──────────────────────────────────────────────────────────────


def _extract_pdf(stream: bytes, mime_type: str) -> ExtractedContent:
    try:
        from pypdf import PdfReader  # type: ignore[import-untyped]
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedMimeType("PDF extraction requires pypdf; add it to pyproject dependencies") from exc

    reader = PdfReader(io.BytesIO(stream))
    if len(reader.pages) > MAX_PDF_PAGES:
        raise UnsupportedMimeType(f"PDF exceeds the {MAX_PDF_PAGES}-page extraction limit")
    spans: list[ExtractedSpan] = []
    low_text_pages: list[int] = []
    for idx, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        # enterprise-gate: broad-except-ok reason=pdf-page-extract-failure-skips-page-not-full-document
        except Exception:
            text = ""
        text = text.strip()
        if len(text) >= 40:
            spans.append(ExtractedSpan(text=text, page=idx))
        else:
            low_text_pages.append(idx)
    ocr_confidences: list[float] = []
    if low_text_pages:
        ocr_spans, ocr_confidences = _ocr_pdf_pages(stream, low_text_pages)
        spans.extend(ocr_spans)
        spans.sort(key=lambda span: span.page or 0)
    return ExtractedContent(
        spans=spans,
        mime_type=mime_type,
        extraction_method="pypdf+ocr" if low_text_pages else "pypdf",
        total_chars=sum(len(s.text) for s in spans),
        extra={
            "page_count": len(reader.pages),
            "ocr_pages": low_text_pages,
            "ocr_mean_confidence": (round(sum(ocr_confidences) / len(ocr_confidences), 2) if ocr_confidences else None),
        },
    )


# ── Word ─────────────────────────────────────────────────────────────


def _extract_docx(stream: bytes, mime_type: str) -> ExtractedContent:
    try:
        from docx import Document  # type: ignore[import-untyped]
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedMimeType("DOCX extraction requires python-docx; ensure it's installed") from exc

    _validate_zip_container(stream)
    doc = Document(io.BytesIO(stream))
    spans: list[ExtractedSpan] = []
    for idx, paragraph in enumerate(doc.paragraphs, start=1):
        text = (paragraph.text or "").strip()
        if text:
            # DOCX has no page concept until reflow. Use paragraph index
            # as provenance so retrieval can still point operators at
            # "paragraph 42".
            spans.append(ExtractedSpan(text=text, page=None, cell_range=f"para {idx}"))
    # Also pull tables — often the most information-dense part of
    # business documents.
    for t_idx, table in enumerate(getattr(doc, "tables", []), start=1):
        for r_idx, row in enumerate(table.rows, start=1):
            cells = [cell.text.strip() for cell in row.cells]
            joined = " | ".join(c for c in cells if c)
            if joined:
                spans.append(
                    ExtractedSpan(
                        text=joined,
                        cell_range=f"table {t_idx} row {r_idx}",
                    )
                )
    return ExtractedContent(
        spans=spans,
        mime_type=mime_type,
        extraction_method="python-docx",
        total_chars=sum(len(s.text) for s in spans),
    )


# ── Excel ────────────────────────────────────────────────────────────


def _extract_xlsx(stream: bytes, mime_type: str) -> ExtractedContent:
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedMimeType("XLSX extraction requires openpyxl; ensure it's installed") from exc

    _validate_zip_container(stream)
    wb = load_workbook(io.BytesIO(stream), data_only=True, read_only=True)
    spans: list[ExtractedSpan] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell).strip() for cell in row if cell is not None and str(cell).strip())
            if not row_text:
                continue
            spans.append(
                ExtractedSpan(
                    text=row_text,
                    sheet=sheet_name,
                    cell_range=None,
                )
            )
    return ExtractedContent(
        spans=spans,
        mime_type=mime_type,
        extraction_method="openpyxl",
        total_chars=sum(len(s.text) for s in spans),
        extra={"sheet_count": len(wb.sheetnames)},
    )


def _extract_pptx(stream: bytes, mime_type: str) -> ExtractedContent:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedMimeType("PPTX extraction requires python-pptx") from exc
    _validate_zip_container(stream)
    presentation = Presentation(io.BytesIO(stream))
    spans: list[ExtractedSpan] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                parts.append(text)
            table = getattr(shape, "table", None) if getattr(shape, "has_table", False) else None
            if table:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        parts.append(row_text)
        joined = _normalise_text("\n".join(parts))
        if joined:
            spans.append(ExtractedSpan(text=joined, page=index, cell_range=f"slide {index}"))
    return ExtractedContent(
        spans=spans,
        mime_type=mime_type,
        extraction_method="python-pptx",
        total_chars=sum(len(span.text) for span in spans),
        extra={"slide_count": len(presentation.slides)},
    )


# ── OCR and conversion helpers ───────────────────────────────────────


def _ocr_image(image: Any) -> tuple[str, float | None]:
    try:
        import pytesseract  # type: ignore[import-untyped]
        from PIL import ImageOps
        from pytesseract import Output  # type: ignore[import-untyped]
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedMimeType("OCR requires Pillow and pytesseract") from exc
    if not shutil.which("tesseract"):
        raise UnsupportedMimeType("OCR requires the Tesseract system binary")
    prepared = ImageOps.autocontrast(ImageOps.grayscale(image))
    if prepared.width < 1800:
        ratio = 1800 / max(prepared.width, 1)
        prepared = prepared.resize((1800, max(1, int(prepared.height * ratio))))
    detected_script = ""
    try:
        orientation = pytesseract.image_to_osd(prepared, timeout=30)
        rotation_line = next(
            (line for line in orientation.splitlines() if line.startswith("Rotate:")),
            "Rotate: 0",
        )
        rotation = int(rotation_line.split(":", 1)[1].strip())
        script_line = next(
            (line for line in orientation.splitlines() if line.startswith("Script:")),
            "Script: Unknown",
        )
        detected_script = script_line.split(":", 1)[1].strip()
        if rotation:
            prepared = prepared.rotate(-rotation, expand=True)
    except (RuntimeError, ValueError):
        pass
    language = _ocr_language_for_script(
        detected_script,
        installed=set(pytesseract.get_languages(config="")),
    )
    try:
        data = pytesseract.image_to_data(
            prepared,
            lang=language,
            config="--oem 1 --psm 6",
            output_type=Output.DICT,
            timeout=120,
        )
    except (RuntimeError, OSError) as exc:
        raise UnsupportedMimeType(f"OCR failed: {exc}") from exc
    words: list[str] = []
    confidences: list[float] = []
    for text, confidence in zip(data.get("text", []), data.get("conf", []), strict=False):
        word = str(text).strip()
        if not word:
            continue
        words.append(word)
        try:
            score = float(confidence)
        except (TypeError, ValueError):
            continue
        if score >= 0:
            confidences.append(score)
    text = _normalise_text(" ".join(words))
    mean = sum(confidences) / len(confidences) if confidences else None
    return text, mean


def _ocr_language_for_script(script: str, *, installed: set[str]) -> str:
    """Choose accurate installed India-first OCR models for an OSD script."""
    configured = os.getenv("AGENTICORG_OCR_LANGUAGES", "").strip()
    requested = (
        configured.split("+")
        if configured
        else {
            "Devanagari": ["eng", "hin", "mar"],
            "Bengali": ["eng", "ben"],
            "Gujarati": ["eng", "guj"],
            "Gurmukhi": ["eng", "pan"],
            "Kannada": ["eng", "kan"],
            "Malayalam": ["eng", "mal"],
            "Tamil": ["eng", "tam"],
            "Telugu": ["eng", "tel"],
            "Arabic": ["eng", "urd"],
        }.get(script, ["eng", "hin"])
    )
    available = [language for language in requested if language in installed]
    if not available:
        raise UnsupportedMimeType("OCR language models are unavailable; install at least tesseract-ocr-eng")
    return "+".join(available)


def _ocr_pdf_pages(stream: bytes, page_numbers: list[int]) -> tuple[list[ExtractedSpan], list[float]]:
    try:
        from pdf2image import convert_from_bytes  # type: ignore[import-untyped]
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedMimeType("Scanned PDF OCR requires pdf2image") from exc
    if len(page_numbers) > MAX_OCR_PAGES:
        raise UnsupportedMimeType(
            f"Scanned PDF requires OCR on {len(page_numbers)} pages; the safe per-document limit is {MAX_OCR_PAGES}"
        )
    spans: list[ExtractedSpan] = []
    confidences: list[float] = []
    for page_number in page_numbers:
        try:
            images = convert_from_bytes(
                stream,
                dpi=300,
                first_page=page_number,
                last_page=page_number,
                fmt="png",
                grayscale=True,
                thread_count=1,
                timeout=120,
            )
        except (OSError, RuntimeError) as exc:
            raise UnsupportedMimeType(f"Scanned PDF OCR failed: {exc}") from exc
        if not images:
            continue
        text, confidence = _ocr_image(images[0])
        if text:
            spans.append(ExtractedSpan(text=text, page=page_number))
        if confidence is not None:
            confidences.append(confidence)
    return spans, confidences


def _extract_image_ocr(stream: bytes, mime_type: str) -> ExtractedContent:
    try:
        from PIL import Image, ImageSequence, UnidentifiedImageError
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedMimeType("Image OCR requires Pillow") from exc
    try:
        image = Image.open(io.BytesIO(stream))
    except (Image.DecompressionBombError, UnidentifiedImageError, OSError) as exc:
        raise UnsupportedMimeType("Image document could not be decoded safely") from exc
    frame_count = int(getattr(image, "n_frames", 1))
    if frame_count > MAX_OCR_PAGES:
        raise UnsupportedMimeType(f"Image document has {frame_count} pages; the safe OCR limit is {MAX_OCR_PAGES}")
    spans: list[ExtractedSpan] = []
    confidences: list[float] = []
    for page, frame in enumerate(ImageSequence.Iterator(image), start=1):
        if frame.width * frame.height > MAX_IMAGE_PIXELS:
            raise UnsupportedMimeType(f"Image page {page} exceeds the {MAX_IMAGE_PIXELS}-pixel OCR limit")
        text, confidence = _ocr_image(frame.copy())
        if text:
            spans.append(ExtractedSpan(text=text, page=page))
        if confidence is not None:
            confidences.append(confidence)
    return ExtractedContent(
        spans=spans,
        mime_type=mime_type,
        extraction_method="tesseract-ocr",
        total_chars=sum(len(span.text) for span in spans),
        extra={
            "page_count": frame_count,
            "ocr_pages": [span.page for span in spans],
            "ocr_mean_confidence": (round(sum(confidences) / len(confidences), 2) if confidences else None),
        },
    )


def _extract_rtf(stream: bytes, mime_type: str) -> ExtractedContent:
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore[import-untyped]
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedMimeType("RTF extraction requires striprtf") from exc
    text = _normalise_text(rtf_to_text(stream.decode("latin-1", errors="replace")))
    return ExtractedContent(
        spans=[ExtractedSpan(text=text)] if text else [],
        mime_type=mime_type,
        extraction_method="striprtf",
        total_chars=len(text),
    )


def _extract_legacy_office(stream: bytes, suffix: str, mime_type: str) -> ExtractedContent:
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if not executable:
        raise UnsupportedMimeType("Legacy/ODF Office extraction requires LibreOffice")
    targets = {
        "doc": ("docx", _extract_docx),
        "odt": ("docx", _extract_docx),
        "xls": ("xlsx", _extract_xlsx),
        "ods": ("xlsx", _extract_xlsx),
        "ppt": ("pptx", _extract_pptx),
        "odp": ("pptx", _extract_pptx),
    }
    target_suffix, extractor = targets[suffix]
    with tempfile.TemporaryDirectory(prefix="agenticorg-office-") as directory:
        source = Path(directory) / f"source.{suffix}"
        source.write_bytes(stream)
        completed = subprocess.run(  # noqa: S603 - executable is resolved from trusted PATH names
            [
                executable,
                "--headless",
                "--nologo",
                "--nolockcheck",
                "--nodefault",
                "--nofirststartwizard",
                "--convert-to",
                target_suffix,
                "--outdir",
                directory,
                str(source),
            ],
            capture_output=True,
            text=True,
            timeout=60,
            check=False,
        )
        converted = Path(directory) / f"source.{target_suffix}"
        if completed.returncode != 0 or not converted.exists():
            raise UnsupportedMimeType("LibreOffice could not convert the uploaded document")
        result = extractor(converted.read_bytes(), mime_type)
        result.extraction_method = f"libreoffice+{result.extraction_method}"
        return result


def _extract_audio(stream: bytes, mime_type: str) -> ExtractedContent:
    raise UnsupportedMimeType(
        "Audio transcription requires ffmpeg + whisper. Enable the "
        "feature flag AGENTICORG_RAG_AUDIO_ENABLED after the deploy "
        "image ships both dependencies and allocates compute budget."
    )


def _extract_video(stream: bytes, mime_type: str) -> ExtractedContent:
    raise UnsupportedMimeType(
        "Video extraction requires ffmpeg + whisper. Enable the feature "
        "flag AGENTICORG_RAG_VIDEO_ENABLED after the deploy image ships "
        "the necessary binaries."
    )


# ── Dispatcher ───────────────────────────────────────────────────────


def extract(stream: bytes, mime_type: str, filename: str = "") -> ExtractedContent:
    """Route ``(stream, mime_type)`` to the correct extractor.

    Unknown MIME types fall through to the filename suffix as a hint;
    truly unrecognised bodies raise ``UnsupportedMimeType`` so the API
    boundary can 415 cleanly.
    """
    mt = (mime_type or "").lower().strip()
    suffix = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    if mt in _TEXT_LIKE_MIMETYPES or suffix in {
        "txt",
        "md",
        "markdown",
        "csv",
        "tsv",
        "json",
        "jsonl",
        "yaml",
        "yml",
        "xml",
        "log",
        "html",
        "htm",
    }:
        if mt == "text/csv" or suffix in {"csv", "tsv"}:
            return _extract_csv(stream, mt)
        if mt in ("application/json", "application/jsonl") or suffix in {"json", "jsonl"}:
            return _extract_json(stream, mt)
        if mt in {"text/html", "application/xhtml+xml"} or suffix in {"html", "htm"}:
            return _extract_html(stream, mt)
        return _extract_plaintext(stream, mt)
    if mt == "message/rfc822" or suffix == "eml":
        return _extract_eml(stream, mt or "message/rfc822")
    if mt == "application/pdf" or suffix == "pdf":
        return _extract_pdf(stream, mt or "application/pdf")
    if mt in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",) or suffix == "docx":
        return _extract_docx(stream, mt or "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    if mt in ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",) or suffix == "xlsx":
        return _extract_xlsx(stream, mt or "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    if mt == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or suffix == "pptx":
        return _extract_pptx(stream, mt or "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    if mt == "application/rtf" or suffix == "rtf":
        return _extract_rtf(stream, mt or "application/rtf")
    if suffix in {"doc", "xls", "ppt", "odt", "ods", "odp"}:
        return _extract_legacy_office(stream, suffix, mt or "application/octet-stream")
    if mt.startswith("image/") or suffix in ("png", "jpg", "jpeg", "webp", "tif", "tiff", "bmp"):
        return _extract_image_ocr(stream, mt or f"image/{suffix}")
    if mt.startswith("audio/") or suffix in ("mp3", "wav", "ogg", "m4a", "flac"):
        return _extract_audio(stream, mt or f"audio/{suffix}")
    if mt.startswith("video/") or suffix in ("mp4", "mov", "mkv", "webm"):
        return _extract_video(stream, mt or f"video/{suffix}")

    # Fall-through: try UTF-8 decode on unknown bodies. If that works and
    # produces meaningful content, index as text — otherwise refuse.
    try:
        text = stream.decode("utf-8")
        if text.strip():
            return ExtractedContent(
                spans=[ExtractedSpan(text=text)],
                mime_type=mt or "application/octet-stream",
                extraction_method="text-fallback",
                total_chars=len(text),
            )
    except UnicodeDecodeError:
        pass
    raise UnsupportedMimeType(f"No extractor registered for mime_type={mt!r} and filename suffix={suffix!r}")
