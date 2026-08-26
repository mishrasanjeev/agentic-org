"""Canonical document extraction and OCR regression coverage."""

from __future__ import annotations

import io
from email.message import EmailMessage

import pytest
from PIL import Image, ImageDraw
from pptx import Presentation
from reportlab.pdfgen import canvas

from core.rag.extractors import ExtractedSpan, UnsupportedMimeType, _ocr_language_for_script, extract


def test_html_and_email_are_extracted_without_active_content() -> None:
    html = b"<h1>Policy</h1><script>steal()</script><p>Refunds need approval.</p>"
    html_result = extract(html, "application/octet-stream", "policy.html")
    assert "Refunds need approval" in html_result.full_text()
    assert "steal" not in html_result.full_text()

    message = EmailMessage()
    message["Subject"] = "Merchant policy"
    message["From"] = "merchant@example.test"
    message["To"] = "agent@example.test"
    message.set_content("Returns are accepted within seven days.")
    email_result = extract(message.as_bytes(), "message/rfc822", "policy.eml")
    assert "Merchant policy" in email_result.full_text()
    assert "seven days" in email_result.full_text()


def test_json_lines_is_advertised_and_malformed_images_fail_cleanly() -> None:
    import asyncio

    from api.v1.knowledge import supported_document_types

    assert ".jsonl" in asyncio.run(supported_document_types())["extensions"]
    with pytest.raises(UnsupportedMimeType, match="decoded safely"):
        extract(b"not-an-image", "image/png", "scan.png")


def test_pptx_has_slide_provenance() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly catalogue"
    slide.placeholders[1].text = "Inventory is refreshed every fifteen minutes."
    body = io.BytesIO()
    presentation.save(body)

    result = extract(
        body.getvalue(),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "catalogue.pptx",
    )

    assert result.extraction_method == "python-pptx"
    assert result.spans[0].page == 1
    assert "Inventory is refreshed" in result.full_text()


def test_scanned_pdf_uses_ocr_for_low_text_pages(monkeypatch) -> None:
    body = io.BytesIO()
    pdf = canvas.Canvas(body)
    pdf.showPage()
    pdf.save()

    monkeypatch.setattr(
        "core.rag.extractors._ocr_pdf_pages",
        lambda _stream, pages: ([ExtractedSpan(text="OCR merchant invoice", page=pages[0])], [98.5]),
    )
    result = extract(body.getvalue(), "application/pdf", "scan.pdf")

    assert result.extraction_method == "pypdf+ocr"
    assert result.extra["ocr_pages"] == [1]
    assert result.extra["ocr_mean_confidence"] == 98.5
    assert "OCR merchant invoice" in result.full_text()


def test_image_document_runs_real_ocr_when_tesseract_is_available() -> None:
    import shutil

    if not shutil.which("tesseract"):
        pytest.skip("Tesseract is exercised in the production-image E2E test")
    image = Image.new("RGB", (1800, 500), "white")
    ImageDraw.Draw(image).text((80, 180), "AGENTICORG MERCHANT CATALOGUE 2026", fill="black")
    body = io.BytesIO()
    image.save(body, format="PNG")

    result = extract(body.getvalue(), "image/png", "scan.png")

    assert result.extraction_method == "tesseract-ocr"
    assert "AGENTICORG" in result.full_text().upper()
    assert result.extra["ocr_pages"] == [1]


def test_supported_types_are_real_and_audio_video_are_truthfully_excluded() -> None:
    import asyncio

    from api.v1.knowledge import supported_document_types

    payload = asyncio.run(supported_document_types())
    for extension in (".pdf", ".docx", ".xlsx", ".pptx", ".png", ".eml", ".odt"):
        assert extension in payload["extensions"]
    assert payload["audio_video_supported"] is False
    assert payload["max_ocr_pages"] > 0


def test_indian_ocr_language_selection_uses_detected_script(monkeypatch) -> None:
    monkeypatch.delenv("AGENTICORG_OCR_LANGUAGES", raising=False)
    installed = {"eng", "hin", "mar", "ben", "tam"}

    assert _ocr_language_for_script("Devanagari", installed=installed) == "eng+hin+mar"
    assert _ocr_language_for_script("Bengali", installed=installed) == "eng+ben"
    assert _ocr_language_for_script("Tamil", installed=installed) == "eng+tam"


def test_generic_mime_tif_and_bmp_dispatch_to_ocr(monkeypatch) -> None:
    monkeypatch.setattr(
        "core.rag.extractors._ocr_image",
        lambda _image: ("scanned merchant policy", 99.0),
    )
    for filename, image_format in (("scan.tif", "TIFF"), ("scan.bmp", "BMP")):
        body = io.BytesIO()
        Image.new("RGB", (100, 100), "white").save(body, format=image_format)
        result = extract(body.getvalue(), "application/octet-stream", filename)
        assert result.extraction_method == "tesseract-ocr"
        assert result.full_text() == "scanned merchant policy"
